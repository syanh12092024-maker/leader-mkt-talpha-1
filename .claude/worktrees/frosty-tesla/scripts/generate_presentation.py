from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_20_slide_presentation():
    prs = Presentation()

    # --- BRAND COLORS ---
    INDIGO = RGBColor(26, 35, 126)
    PURPLE = RGBColor(142, 68, 173)
    CYAN = RGBColor(0, 188, 212)
    WHITE = RGBColor(255, 255, 255)
    GRAY_LIGHT = RGBColor(236, 240, 241)

    # --- IMAGE PATHS ---
    base_path = "/Users/macbook/.gemini/antigravity/brain/0a7ff803-e424-4384-bea7-509a67229fa2/"
    hero_img = base_path + "vibe_coding_hero_1774855313381.png"
    flow_img = base_path + "automation_flow_infographic_1774855331442.png"
    roadmap_img = base_path + "future_roadmap_vision_1774857121022.png"
    omni_img = base_path + "omnichannel_matrix_1774857138102.png"
    badge_img = base_path + "ai_innovation_badge_1774857153376.png"

    def add_image_slide(prs, img_path, title_text, subtitle_text=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, 0, 0, height=prs.slide_height)
        
        # Add Overlay
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(3))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = title_text
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.size = Pt(24)
            p2.font.color.rgb = GRAY_LIGHT
            p2.alignment = PP_ALIGN.CENTER
        return slide

    def add_bullet_slide(prs, title, points, heading=None):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        
        if heading:
            p = tf.add_paragraph()
            p.text = heading
            p.font.bold = True
            p.font.size = Pt(24)
            p.space_after = Pt(10)
            
        for p_title, p_detail in points:
            p = tf.add_paragraph()
            p.text = p_title
            p.font.bold = True
            p.font.size = Pt(20)
            p.font.color.rgb = INDIGO
            
            p2 = tf.add_paragraph()
            p2.text = f"  → {p_detail}"
            p2.font.size = Pt(16)
            p2.space_after = Pt(8)
        return slide

    # --- SECTION 1: CONTEXT ---
    add_image_slide(prs, hero_img, "Talpha Broadcast", "Vibe Coding Challenge 2026: Kiến tạo tương lai bằng AI")
    
    add_bullet_slide(prs, "Level Up's Digital DNA", [
        ("AI First Strategy", "Tích hợp AI vào mọi quy trình vận hành cốt lõi của công ty."),
        ("Tối ưu hóa nguồn lực", "Biến nhân sự từ người thực thi thủ công thành người quản trị AI."),
        ("Đột phá trải nghiệm", "Cung cấp sự phản hồi và chăm sóc tức thì cho hàng triệu khách hàng.")
    ])
    
    add_bullet_slide(prs, "The 'Vibe Coding' Paradigm", [
        ("Idea to Impact", " AI rút ngắn khoảng cách từ ý tưởng đến sản phẩm từ vài tháng xuống vài ngày."),
        ("Human-Centric Design", "Code phục vụ con người, không phải con người phục vụ code."),
        ("Vận hành thông minh", "Tự động hóa các tác vụ lặp đi lặp lại để tập trung vào sáng tạo.")
    ])
    
    add_bullet_slide(prs, "Nội dung thuyết trình", [
        ("Thực trạng & Thách thức", "Rào cản Facebook và chi phí vận hành."),
        ("Hệ năng lực cốt lõi", "Chi tiết chức năng của Talpha Broadcast."),
        ("Thế mạnh kỹ thuật", "Tại sao hệ thống này lại khác biệt?"),
        ("Tầm nhìn dài hạn", "Lộ trình 2026 - 2027.")
    ])

    # --- SECTION 2: PAIN POINTS ---
    add_bullet_slide(prs, "Rào cản Facebook 24h & 7 ngày", [
        ("The 24-Hour Policy", "Chặn tin nhắn miễn phí sau 24h từ lần tương tác cuối."),
        ("The 7-Day Window", "Kết thúc hoàn toàn khả năng nhắn tin chủ động sau 7 ngày."),
        ("Hệ lụy", "Mất kết nối với 80% khách hàng tiềm năng đã từng nhắn tin/mua hàng.")
    ])
    
    add_bullet_slide(prs, "Chi phí của sự thiếu hiệu quả", [
        ("Năng suất thấp", "Nhân viên Telesale/CSKH tốn hàng giờ để lọc và nhắn tin thủ công."),
        ("Sai sót dữ liệu", "Dễ gửi nhầm, gửi lặp hoặc quên khách hàng do quản lý Excel/CRM rời rạc."),
        ("ROI sụt giảm", "Chi phí Marketing tăng nhưng tỷ lệ chuyển đổi từ khách cũ thấp.")
    ])
    
    add_bullet_slide(prs, "Data Disconnect: CRM vs Reality", [
        ("Pancake CRM", "Dữ liệu khách hàng dồi dào nhưng khó khai thác hàng loạt."),
        ("POS Records", "Dữ liệu đơn hàng chưa được tận dụng để tái tiếp thị (Remarketing)."),
        ("Talpha Bridge", "Sứ mệnh kết nối các nền tảng thành một luồng dữ liệu thông suốt.")
    ])
    
    add_bullet_slide(prs, "Cơ hội bị bỏ lỡ", [
        ("Silent Churn", "Khách hàng rời đi lặng lẽ vì không được nhắc nhớ về thương hiệu."),
        ("Re-engagement Gap", "Thiếu công cụ để kích hoạt lại các 'Lead' đã nguội."),
        ("Mục tiêu", "Biến dữ liệu tĩnh thành dòng tiền động.")
    ])

    # --- SECTION 3: FUNCTIONALITY ---
    add_image_slide(prs, flow_img, "Core Architecture", "Hệ thống kết nối đa nền tảng")
    
    add_bullet_slide(prs, "Audience Filtering thông minh", [
        ("Shop & Page Selection", "Quản lý tập trung hàng chục fanpage trên nhiều quốc gia (Saudi, UAE, Japan...)."),
        ("Purchase Status", "Lọc khách hàng đã mua, chưa mua, hoặc khách hàng VIP."),
        ("Time-based Trigger", "Nhắm mục tiêu khách hàng theo thời gian tương tác cuối cùng.")
    ])
    
    add_bullet_slide(prs, "The Re-engagement Engine", [
        ("Human Agent Automation", "Tự động áp dụng tag HUMAN_AGENT để bắn tin trong khung 7 ngày."),
        ("Contextual Messaging", "Soạn kịch bản cá nhân hóa cho từng nhóm đối tượng khách hàng."),
        ("Template Central", "Quản lý kịch bản gửi tin chuyên nghiệp và đồng bộ.")
    ])
    
    add_bullet_slide(prs, "Multimedia Cloud Asset", [
        ("Auto-Upload Engine", "Tự động đẩy hình ảnh lên Cloud để lấy link API tương thích."),
        ("High-Speed Delivery", "Đảm bảo hình ảnh hiện thị tức thì khi khách hàng mở tin nhắn."),
        ("Asset Security", "Bảo mật và tối ưu hóa tài nguyên hình ảnh.")
    ])

    # --- SECTION 4: STRENGTHS ---
    add_bullet_slide(prs, "High-Performance Persistence", [
        ("Smart Pagination", "Quét hàng chục ngàn hội thoại mà không bị treo hệ thống."),
        ("Server-side Dedup", "Cơ chế chống gửi lặp (5 phút safety window) bảo vệ Fanpage."),
        ("Data Integrity", "Đảm bảo mọi khách hàng đều được xử lý đúng kịch bản.")
    ])
    
    add_bullet_slide(prs, "Self-Healing Hybrid API", [
        ("Adaptive Routing", "Tự động chuyển đổi giữa Pancake API và FB Graph API."),
        ("Token Management", "Quản lý Page Access Token tự động và bảo mật."),
        ("Resilience", "Hệ thống tự phục hồi khi gặp lỗi kết nối API.")
    ])
    
    add_bullet_slide(prs, "Scalability: Sẵn sàng cho 10,000+ KH", [
        ("Concurrency Control", "Xử lý hàng loạt tin nhắn mà không vi phạm Spam policy."),
        ("Lightweight Backend", "Vận hành mượt mà trên hạ tầng Serverless (Vercel)."),
        ("Ready for Scale", "Dễ dàng mở rộng ra hàng trăm nhân sự và hàng triệu khách hàng.")
    ])

    # --- SECTION 5: UI & ANALYTICS ---
    add_bullet_slide(prs, "Dashboard: Power in Simplicity", [
        ("Real-time Tracking", "Theo dõi chính xác số tin nhắn đã gửi, thành công, và lỗi."),
        ("User Experience", "Giao diện tối giản, chỉ mất 1 phút để thiết lập một chiến dịch."),
        ("Visibility", "Công khai minh bạch tiến độ bắn tin cho toàn bộ team.")
    ])
    
    add_bullet_slide(prs, "Error Intelligence", [
        ("Error Code Translation", "Dịch các mã lỗi kỹ thuật (#10, #551) sang ngôn ngữ dễ hiểu."),
        ("Direct Action", "Gợi ý cách sửa lỗi ngay khi hệ thống phát hiện vấn đề."),
        ("Detailed Logs", "Lưu trữ lịch sử gửi tin để phân tích hiệu quả.")
    ])

    # --- SECTION 6: FUTURE ---
    add_bullet_slide(prs, "Quantifiable ROI (Dự kiến)", [
        ("Tốc độ vận hành", "Tăng 900% so với làm thủ công."),
        ("Tỷ lệ chuyển đổi", "Tăng 15% doanh số từ tệp khách hàng cũ."),
        ("Tiết kiệm chi phí", "Giảm 50% áp lực nhân sự CSKH trong các mùa cao điểm.")
    ])
    
    add_image_slide(prs, roadmap_img, "Vision 2026 - 2027", "Lộ trình Omni-Channel & AI Agent")
    
    add_bullet_slide(prs, "Roadmap & Thế hệ kế tiếp", [
        ("Omni-Channel Extension", "Mở rộng sang WhatsApp, Zalo và Website Chat."),
        ("AI Chat Agents", "Tích hợp Llama 3 để AI tự động tư vấn và chốt đơn sau khi bắn tin."),
        ("Predictive Sales", "Dùng AI dự báo thời điểm khách hàng sắp hết hàng để chủ động bắn tin.")
    ])
    
    add_image_slide(prs, badge_img, "Idea to Impact", "Sẵn sàng vươn tầm cao mới")

    # --- SAVE ---
    output_path = "/Users/macbook/Desktop/Talpha_Broadcast_Comprehensive_VibeCoding2026.pptx"
    prs.save(output_path)
    print(f"20-slide presentation saved to: {output_path}")

if __name__ == "__main__":
    create_20_slide_presentation()
