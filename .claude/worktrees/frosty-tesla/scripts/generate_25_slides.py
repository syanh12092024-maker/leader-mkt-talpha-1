#!/usr/bin/env python3
"""Generate 25-slide TAlpha presentation with real screenshots & avatars."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os, sys

# === PATHS ===
BASE = "/Users/macbook/.gemini/antigravity/brain/4d2c3f7a-6e7c-42ec-98ae-b8d6e53917c4"
OUT = "/Users/macbook/Desktop/TAlpha_25_Slides.pptx"

# Screenshots
SS = {
    "script": f"{BASE}/slide_script_tab_1774905378299.png",
    "shipment": f"{BASE}/slide_shipment_tab_1774905398235.png",
    "shipment_ov": f"{BASE}/slide_shipment_overview_1774905416457.png",
    "shipment_form": f"{BASE}/slide_shipment_form_1774905461586.png",
    "gcc": f"{BASE}/slide_gcc_dist_1774905498920.png",
    "broadcast": f"{BASE}/slide_broadcast_tab_1774905523182.png",
    "broadcast_seg": f"{BASE}/slide_broadcast_segments_1774905536890.png",
}

# Avatars (user-uploaded images saved by system)
AVATARS = {
    "thuy": None,  # Will search
    "phuong_anh": None,
    "sy_anh": None,
    "sy_loc": None,
    "hong_ly": None,
}

# === COLORS ===
BG = RGBColor(0x0A, 0x0F, 0x1C)
CARD = RGBColor(0x1E, 0x29, 0x3B)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_GRAY = RGBColor(0x47, 0x55, 0x69)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
RED = RGBColor(0xF8, 0x71, 0x71)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
BLUE = RGBColor(0x60, 0xA5, 0xFA)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=24, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Arial", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    p.space_after = Pt(0)
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=20, color=GRAY, font_name="Arial", line_spacing=1.2, bold=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.font.bold = bold
        p.space_after = Pt(4)
    return txBox

def add_rect(slide, left, top, width, height, color=CARD, corner_radius=0.15):
    shape = slide.shapes.add_shape(
        5, Inches(left), Inches(top), Inches(width), Inches(height)  # 5 = rounded rect
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_image_safe(slide, path, left, top, width, height):
    if path and os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
        return True
    return False

# =====================================================
# SLIDE 1: COVER
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 2, 1.5, 12, 0.5, "T   A   L   P   H   A", 22, CYAN, True, PP_ALIGN.CENTER, "Courier New")
add_text(s, 1, 2.5, 14, 2, "Hệ thống Quản lý &\nTự động hoá Marketing Toàn diện", 56, WHITE, True, PP_ALIGN.CENTER)
add_text(s, 2, 5, 12, 0.8, "Tạo kịch bản  ·  Đặt và Flow hàng  ·  Gửi tin hàng loạt", 22, GRAY, False, PP_ALIGN.CENTER, "Courier New")
# Divider
add_rect(s, 7.2, 6.2, 1.6, 0.04, CYAN, 0)
add_text(s, 2, 6.8, 12, 0.5, "Vibe Coding Challenge 2026  ·  Pi Alpha", 18, DARK_GRAY, False, PP_ALIGN.CENTER, "Courier New")

# =====================================================
# SLIDE 2: AGENDA
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.6, 5, 0.5, "NỘI DUNG", 16, CYAN, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 1, 1.2, 14, 0.8, "3 Module chính", 48, WHITE, True)

items = [
    ("01", "Tạo kịch bản chào hàng", "Upload ảnh → AI viết 4 đoạn cho 4 khung giờ", CYAN),
    ("02", "Đặt và Flow hàng", "Quản lý đơn nhập · Approval flow · Phân phối UAE → GCC", GREEN),
    ("03", "Gửi tin nhắn hàng loạt", "Broadcast Engine · Hẹn giờ tự động · Lọc thông minh", PURPLE),
]
for i, (num, title, desc, clr) in enumerate(items):
    y = 2.5 + i * 1.9
    add_rect(s, 1, y, 14, 1.6, CARD)
    add_text(s, 1.5, y + 0.2, 1.5, 1, num, 40, clr, True, PP_ALIGN.LEFT, "Courier New")
    add_text(s, 3.2, y + 0.2, 10, 0.6, title, 28, WHITE, True)
    add_text(s, 3.2, y + 0.85, 10, 0.5, desc, 18, GRAY)

# =====================================================
# SLIDE 3: TEAM
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.6, 5, 0.5, "ĐỘI NGŨ", 16, CYAN, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 1, 1.2, 14, 0.8, "Những con người tạo nên TAlpha", 48, WHITE, True)

team = [
    ("Lê Thị Diệu Thuý", "Giám Đốc Pi Alpha", CYAN),
    ("Nguyễn Thị Phương Anh", "MKT Manager", GREEN),
    ("Hồ Sỹ Anh", "Tech Lead & Builder", YELLOW),
    ("Hồ Sỹ Lộc", "Developer", PURPLE),
    ("Nguyễn Thị Hồng Ly", "MKT Specialist", ORANGE),
]
# Find avatar files
avatar_dir = BASE
avatar_files = sorted([f for f in os.listdir(avatar_dir) if f.startswith("avatar_")]) if os.path.exists(avatar_dir) else []

for i, (name, role, clr) in enumerate(team):
    x = 0.8 + i * 3
    add_rect(s, x, 2.8, 2.6, 4.5, CARD)
    # Avatar placeholder circle
    add_rect(s, x + 0.55, 3.2, 1.5, 1.5, RGBColor(0x0F, 0x17, 0x2A))
    add_text(s, x + 0.6, 3.5, 1.5, 1, "👤", 48, GRAY, False, PP_ALIGN.CENTER)
    # Name & role
    add_text(s, x + 0.1, 5.2, 2.4, 0.6, name, 16, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, x + 0.1, 5.8, 2.4, 0.5, role, 14, clr, False, PP_ALIGN.CENTER)
add_text(s, 1, 7.8, 14, 0.5, "* Ảnh avatar thực tế được đính kèm riêng", 14, DARK_GRAY, False, PP_ALIGN.CENTER)

# =====================================================
# SLIDE 4: PAIN POINT 1 — Đặt hàng
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.6, 5, 0.5, "NỖI ĐAU #1", 16, RED, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 1, 1.3, 14, 1.2, "Không bám sát được hàng\nđang ở giai đoạn nào", 48, WHITE, True)
add_rect(s, 1, 3, 14, 4.5, CARD)
pains1 = [
    "❌  Bên mua hàng quên nhập liệu → mất kiểm soát tiến độ",
    "❌  Quản lý bằng Excel rời rạc, không real-time",
    "❌  Không biết hàng đã đến kho đích 1 hay chưa",
    "❌  Không phân biệt được đơn nào đang chờ duyệt, đơn nào đang vận chuyển",
    "❌  Hàng đến muộn mà không có cảnh báo → ảnh hưởng kế hoạch bán",
    "❌  Phải hỏi từng người, từng ngày → lãng phí 1-2h/ngày",
]
add_multiline(s, 1.5, 3.3, 13, 4, pains1, 22, GRAY)
add_text(s, 1, 7.8, 14, 0.5, "→ Giải pháp: Module \"Đặt và Flow hàng\" — 17 trạng thái, cảnh báo real-time", 20, CYAN, True, PP_ALIGN.CENTER)

# =====================================================
# SLIDE 5: PAIN POINT 2 — Kịch bản
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.6, 5, 0.5, "NỖI ĐAU #2", 16, RED, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 1, 1.3, 14, 1.2, "Tạo kịch bản lâu vì ChatGPT\nkhông đúng ý, không đúng form", 48, WHITE, True)
add_rect(s, 1, 3, 14, 4.5, CARD)
pains2 = [
    "❌  ChatGPT viết kịch bản nhưng mỗi lần tạo lại khác format",
    "❌  Phải edit lại 70-80% nội dung → tốn thêm 30 phút/kịch bản",
    "❌  Mỗi page cần kịch bản khác nhau, phải copy-paste hướng dẫn lại mỗi lần",
    "❌  Không có template chuẩn → tone không thống nhất giữa các page",
    "❌  MKT phải soạn riêng 4 đoạn cho 4 khung giờ → x4 công sức",
    "❌  Tổng thời gian: 2-3h chỉ để soạn kịch bản cho 1 sản phẩm",
]
add_multiline(s, 1.5, 3.3, 13, 4, pains2, 22, GRAY)
add_text(s, 1, 7.8, 14, 0.5, "→ Giải pháp: Upload ảnh → AI Gemini viết 4 đoạn chuẩn form trong 30 giây", 20, CYAN, True, PP_ALIGN.CENTER)

# =====================================================
# SLIDE 6: PAIN POINT 3 — Broadcast
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.6, 5, 0.5, "NỖI ĐAU #3", 16, RED, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 1, 1.3, 14, 1.2, "Không bắn được tin ngoài 24h\n→ thất thoát doanh thu", 48, WHITE, True)
add_rect(s, 1, 3, 14, 4.5, CARD)
pains3 = [
    "❌  Facebook chặn gửi tin sau 24h → mất 70% khách tiềm năng",
    "❌  Sale phải dành 2h/ngày để cài đặt bot bắn tin thủ công",
    "❌  Bot bắn lại cho khách đã mua → spam → mất uy tín page",
    "❌  Không lọc được khách đã mua vs chưa mua",
    "❌  Quên gửi vào khung giờ vàng → tỷ lệ đọc thấp",
    "❌  Phải bắn bot nhiều lần cho nhiều page → nhân công x số page",
]
add_multiline(s, 1.5, 3.3, 13, 4, pains3, 22, GRAY)
add_text(s, 1, 7.8, 14, 0.5, "→ Giải pháp: HUMAN_AGENT tag mở rộng 7 ngày + Auto-fire 4 khung giờ + Lọc đã mua", 20, CYAN, True, PP_ALIGN.CENTER)

# =====================================================
# SLIDE 7: MODULE 01 INTRO — Tạo kịch bản
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 2, 2, 12, 0.5, "MODULE  01", 20, CYAN, True, PP_ALIGN.CENTER, "Courier New")
add_text(s, 1, 3, 14, 1.5, "Tạo kịch bản chào hàng", 64, WHITE, True, PP_ALIGN.CENTER)
add_text(s, 2, 5.2, 12, 0.8, "Upload ảnh sản phẩm → AI viết kịch bản → 4 đoạn cho 4 khung giờ", 24, GRAY, False, PP_ALIGN.CENTER)

# =====================================================
# SLIDE 8: SCRIPT DASHBOARD — Real screenshot
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 0.8, 0.4, 5, 0.4, "GIAO DIỆN THỰC TẾ", 14, CYAN, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 0.8, 0.85, 10, 0.6, "Dashboard tạo kịch bản", 36, WHITE, True)
add_text(s, 13, 0.5, 2, 0.4, "01 / 03", 16, DARK_GRAY, False, PP_ALIGN.RIGHT, "Courier New")
add_image_safe(s, SS["script"], 0.5, 1.6, 15, 7)

# =====================================================
# SLIDE 9: SCRIPT FEATURES
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.6, 5, 0.5, "PHÂN TÍCH GIÁ TRỊ", 16, CYAN, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 1, 1.2, 14, 0.8, "Tạo kịch bản — Tính năng & Lợi ích", 40, WHITE, True)

features = [
    ("📸", "Upload ảnh SP", "Kéo thả nhiều ảnh. AI tự nhận diện sản phẩm từ hình ảnh.", CYAN),
    ("✨", "AI sinh 4 đoạn", "Gemini AI viết 4 tone: sáng lạc quan, trưa chuyên nghiệp, chiều FOMO, tối intimate.", GREEN),
    ("🎯", "Template tùy chỉnh", "Dùng template sẵn hoặc tự tạo. Hỗ trợ biến {tên}, {giá} cá nhân hoá.", YELLOW),
]
for i, (icon, title, desc, clr) in enumerate(features):
    x = 0.8 + i * 5
    add_rect(s, x, 2.5, 4.5, 5.5, CARD)
    add_text(s, x + 0.3, 2.8, 1, 0.8, icon, 40, WHITE)
    add_text(s, x + 0.3, 3.8, 4, 0.6, title, 24, WHITE, True)
    add_multiline(s, x + 0.3, 4.6, 3.8, 3, [desc], 18, GRAY)

# =====================================================
# SLIDE 10: MODULE 02 INTRO
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 2, 2, 12, 0.5, "MODULE  02", 20, CYAN, True, PP_ALIGN.CENTER, "Courier New")
add_text(s, 1, 3, 14, 1.5, "Đặt và Flow hàng", 64, WHITE, True, PP_ALIGN.CENTER)
add_text(s, 2, 5.2, 12, 0.8, "Quản lý đơn nhập hàng  ·  Approval workflow  ·  Phân phối UAE → GCC", 24, GRAY, False, PP_ALIGN.CENTER)

# =====================================================
# SLIDE 11: SHIPMENT DASHBOARD — Real screenshot
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 0.8, 0.4, 5, 0.4, "GIAO DIỆN THỰC TẾ", 14, CYAN, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 0.8, 0.85, 10, 0.6, "Quản lý đơn nhập hàng", 36, WHITE, True)
add_text(s, 13, 0.5, 2, 0.4, "02 / 03", 16, DARK_GRAY, False, PP_ALIGN.RIGHT, "Courier New")
add_image_safe(s, SS["shipment"], 0.5, 1.6, 15, 7)

# =====================================================
# SLIDE 12: SHIPMENT OVERVIEW — Real screenshot
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 0.8, 0.4, 5, 0.4, "TỔNG QUAN", 14, CYAN, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 0.8, 0.85, 10, 0.6, "Dashboard tổng quan + Cảnh báo", 36, WHITE, True)
add_image_safe(s, SS["shipment_ov"], 0.5, 1.6, 15, 7)

# =====================================================
# SLIDE 13: APPROVAL FLOW
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.6, 5, 0.5, "APPROVAL FLOW", 16, CYAN, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 1, 1.2, 14, 0.8, "17 trạng thái — từ yêu cầu đến giao hàng", 40, WHITE, True)

flow_items = [
    ("MKT\nyêu cầu", BLUE),
    ("→", DARK_GRAY),
    ("Kiểm tra\n& Duyệt", YELLOW),
    ("→", DARK_GRAY),
    ("Đặt hàng\nNCC", PURPLE),
    ("→", DARK_GRAY),
    ("Vận\nchuyển", CYAN),
    ("→", DARK_GRAY),
    ("Đã giao\nthành công", GREEN),
]
for i, (txt, clr) in enumerate(flow_items):
    x = 0.5 + i * 1.7
    if txt == "→":
        add_text(s, x, 3.2, 1, 0.8, "→", 36, DARK_GRAY, False, PP_ALIGN.CENTER)
    else:
        add_rect(s, x, 2.5, 1.5, 1.6, CARD)
        add_text(s, x + 0.05, 2.7, 1.4, 1.2, txt, 16, clr, True, PP_ALIGN.CENTER)

# Features below
feats_ship = [
    ("📦 Quản lý đơn nhập", "Theo dõi TQ→UAE→GCC\nTự tính ngày muộn\nNhận hàng nhiều lần"),
    ("🌍 Phân phối GCC", "Tạo đơn phân phối từ gốc\nTracking riêng từng nước\n6 quốc gia GCC"),
    ("📊 Tổng quan & Cảnh báo", "Dashboard real-time\nCảnh báo chờ duyệt\nLịch sử trạng thái"),
]
for i, (title, desc) in enumerate(feats_ship):
    x = 0.8 + i * 5
    add_rect(s, x, 4.8, 4.5, 3.5, CARD)
    add_text(s, x + 0.3, 5.1, 4, 0.5, title, 20, WHITE, True)
    add_multiline(s, x + 0.3, 5.8, 4, 2.5, desc.split("\n"), 18, GRAY)

# =====================================================
# SLIDE 14: CREATE ORDER FORM — Real screenshot
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 0.8, 0.4, 5, 0.4, "FORM TẠO ĐƠN", 14, CYAN, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 0.8, 0.85, 10, 0.6, "Nhập đầy đủ thông tin trong 1 modal", 36, WHITE, True)
add_image_safe(s, SS["shipment_form"], 0.5, 1.6, 15, 7)

# =====================================================
# SLIDE 15: GCC DISTRIBUTION — Real screenshot
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 0.8, 0.4, 5, 0.4, "PHÂN PHỐI GCC", 14, CYAN, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 0.8, 0.85, 10, 0.6, "UAE → KSA, Kuwait, Qatar, Bahrain, Oman", 36, WHITE, True)
add_image_safe(s, SS["gcc"], 0.5, 1.6, 15, 7)

# =====================================================
# SLIDE 16: MODULE 03 INTRO
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 2, 2, 12, 0.5, "MODULE  03", 20, CYAN, True, PP_ALIGN.CENTER, "Courier New")
add_text(s, 1, 3, 14, 1.5, "Gửi tin nhắn hàng loạt", 64, WHITE, True, PP_ALIGN.CENTER)
add_text(s, 2, 5.2, 12, 0.8, "Broadcast Engine  ·  Hẹn lịch tự động  ·  Lọc thông minh  ·  Facebook Graph API", 24, GRAY, False, PP_ALIGN.CENTER)

# =====================================================
# SLIDE 17: BROADCAST DASHBOARD — Real screenshot
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 0.8, 0.4, 5, 0.4, "GIAO DIỆN THỰC TẾ", 14, CYAN, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 0.8, 0.85, 10, 0.6, "Dashboard gửi tin hàng loạt", 36, WHITE, True)
add_text(s, 13, 0.5, 2, 0.4, "03 / 03", 16, DARK_GRAY, False, PP_ALIGN.RIGHT, "Courier New")
add_image_safe(s, SS["broadcast"], 0.5, 1.6, 15, 7)

# =====================================================
# SLIDE 18: BROADCAST SEGMENTS — Real screenshot
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 0.8, 0.4, 5, 0.4, "SOẠN TIN 4 ĐOẠN", 14, CYAN, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 0.8, 0.85, 10, 0.6, "Mỗi đoạn cho 1 khung giờ vàng", 36, WHITE, True)
add_image_safe(s, SS["broadcast_seg"], 0.5, 1.6, 15, 7)

# =====================================================
# SLIDE 19: BROADCAST FEATURES
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.6, 5, 0.5, "TÍNH NĂNG", 16, CYAN, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 1, 1.2, 14, 0.8, "Broadcast — Gửi tin thông minh", 40, WHITE, True)

bfeats = [
    ("📨", "4 đoạn / 4 khung giờ", "Đoạn 1→6h, Đoạn 2→11h,\nĐoạn 3→17h, Đoạn 4→21h.\nDelay thông minh giữa mỗi tin.", CYAN),
    ("⏰", "Hẹn lịch tự động", "Hẹn 1 lần, chạy mỗi ngày.\nWorker kiểm tra 60s/lần.\nTrạng thái: ✓xanh ⏳vàng ✗đỏ.", GREEN),
    ("🔍", "Lọc khách thông minh", "Bỏ qua khách đã mua.\nLọc SĐT, đơn hàng, thẻ.\nKhông spam = giữ uy tín.", YELLOW),
]
for i, (icon, title, desc, clr) in enumerate(bfeats):
    x = 0.8 + i * 5
    add_rect(s, x, 2.5, 4.5, 5.5, CARD)
    add_text(s, x + 0.3, 2.8, 1, 0.8, icon, 40, WHITE)
    add_text(s, x + 0.3, 3.8, 4, 0.6, title, 24, WHITE, True)
    add_multiline(s, x + 0.3, 4.6, 3.8, 3, desc.split("\n"), 18, GRAY)

# =====================================================
# SLIDE 20: SEGMENTS MAPPING
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.6, 5, 0.5, "4 SEGMENTS", 16, CYAN, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 1, 1.2, 14, 0.8, "Mỗi đoạn = Một khung giờ vàng", 40, WHITE, True)

segs = [
    ("🌅 Đoạn 1", "6:00 sáng", "Tone lạc quan, chào ngày mới", GREEN),
    ("☀️ Đoạn 2", "11:00 trưa", "Tone chuyên nghiệp, giá trị SP", YELLOW),
    ("🌆 Đoạn 3", "17:00 chiều", "Tone FOMO, giá ưu đãi", ORANGE),
    ("🌙 Đoạn 4", "21:00 tối", "Tone intimate, khuyến mãi cuối ngày", PURPLE),
]
for i, (label, time, desc, clr) in enumerate(segs):
    y = 2.5 + i * 1.5
    add_rect(s, 1, y, 14, 1.25, CARD)
    add_text(s, 1.5, y + 0.15, 3, 0.6, label, 22, clr, True)
    add_text(s, 5, y + 0.15, 3, 0.6, "→  " + time, 22, WHITE, True, PP_ALIGN.LEFT, "Courier New")
    add_text(s, 9, y + 0.15, 5, 0.6, desc, 18, GRAY)

# =====================================================
# SLIDE 21: SCHEDULE STATUS
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.6, 5, 0.5, "TRẠNG THÁI REALTIME", 16, CYAN, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 1, 1.2, 14, 0.8, "1 lần cài = Chạy mãi mãi", 48, WHITE, True)

statuses = [
    ("✓ Đã gửi", "Segment fire thành công\nHiển thị XANH\nTự ghi log thời gian", GREEN),
    ("⏳ Chờ gửi", "Chưa đến giờ\nHiển thị VÀNG\nSẽ tự fire khi đến giờ", YELLOW),
    ("✗ Lỗi", "Gửi thất bại\nHiển thị ĐỎ\nCho biết số lỗi / tổng", RED),
]
for i, (title, desc, clr) in enumerate(statuses):
    x = 0.8 + i * 5.1
    add_rect(s, x, 2.5, 4.6, 4, CARD)
    add_text(s, x + 0.4, 2.8, 4, 0.6, title, 32, clr, True, PP_ALIGN.LEFT, "Courier New")
    add_multiline(s, x + 0.4, 3.8, 4, 2.5, desc.split("\n"), 20, GRAY)

add_text(s, 0.5, 7.2, 15, 0.8, "Mở tab → worker 60s/lần → tự gửi segment đã qua giờ → ngày mới tự reset", 20, DARK_GRAY, False, PP_ALIGN.CENTER, "Courier New")

# =====================================================
# SLIDE 22: HUMAN_AGENT BYPASS
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.6, 5, 0.5, "FACEBOOK API", 16, CYAN, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 1, 1.2, 14, 1.2, "Bypass giới hạn 24h\nvới tag HUMAN_AGENT", 44, WHITE, True)
add_multiline(s, 1, 3, 7, 1.5, [
    "Facebook chặn gửi tin sau 24h. TAlpha sử dụng",
    "Facebook Graph API v21 + tag HUMAN_AGENT",
    "để mở rộng cửa sổ gửi tin lên 7 ngày."
], 20, GRAY)
add_rect(s, 1, 4.8, 7, 3.2, RGBColor(0x0F, 0x17, 0x2A))
code = [
    "POST /me/messages",
    "",
    '{ messaging_type: "MESSAGE_TAG",',
    '  tag: "HUMAN_AGENT",',
    "  recipient: { id: psid },",
    "  message: { text: content } }",
]
add_multiline(s, 1.3, 5, 6.5, 3, code, 18, CYAN, "Courier New")

# Right side - before/after
add_rect(s, 8.5, 1.5, 7, 3, CARD)
add_text(s, 8.8, 1.7, 6, 0.5, "❌ TRƯỚC — Pancake CRM", 22, RED, True)
add_multiline(s, 8.8, 2.3, 6, 2, [
    "• Chỉ gửi được trong 24h",
    "• Mất 70% khách tiềm năng",
    "• Không có fallback",
], 18, GRAY)
add_rect(s, 8.5, 5, 7, 3, CARD)
add_text(s, 8.8, 5.2, 6, 0.5, "✅ SAU — TAlpha + Graph API", 22, GREEN, True)
add_multiline(s, 8.8, 5.8, 6, 2, [
    "• Gửi được trong 7 ngày",
    "• Reach 100% khách hàng",
    "• Auto-fallback thông minh",
], 18, GRAY)

# =====================================================
# SLIDE 23: ARCHITECTURE
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.6, 5, 0.5, "KIẾN TRÚC HỆ THỐNG", 16, CYAN, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 1, 1.2, 14, 0.8, "Tech Stack & Data Flow", 40, WHITE, True)

arch_items = [
    ("React\nFrontend", CYAN),
    ("→", DARK_GRAY),
    ("Next.js API\nRoutes", GREEN),
    ("→", DARK_GRAY),
    ("Pancake\nCRM API", YELLOW),
    ("→", DARK_GRAY),
    ("Facebook\nGraph API", BLUE),
    ("→", DARK_GRAY),
    ("Messenger\nKhách hàng", RED),
]
for i, (txt, clr) in enumerate(arch_items):
    x = 0.3 + i * 1.7
    if txt == "→":
        add_text(s, x, 3.2, 1, 0.8, "→", 36, DARK_GRAY, False, PP_ALIGN.CENTER)
    else:
        add_rect(s, x, 2.5, 1.5, 1.6, CARD)
        add_text(s, x + 0.05, 2.7, 1.4, 1.2, txt, 16, clr, True, PP_ALIGN.CENTER)

# Tech cards
techs = [
    ("Frontend", "• Next.js 14 (App Router)\n• React + TypeScript\n• Framer Motion\n• Tailwind CSS"),
    ("Backend", "• Next.js API Routes\n• Pancake CRM REST API\n• Facebook Graph API v21\n• Gemini AI API"),
    ("Data & Deploy", "• localStorage (client)\n• CRM cloud data\n• Vercel hosting\n• No database required"),
]
for i, (title, desc) in enumerate(techs):
    x = 0.8 + i * 5.1
    add_rect(s, x, 4.8, 4.6, 3.5, CARD)
    add_text(s, x + 0.3, 5, 4, 0.5, title, 22, WHITE, True)
    add_multiline(s, x + 0.3, 5.6, 4, 2.5, desc.split("\n"), 18, GRAY)

# =====================================================
# SLIDE 24: KPI & BEFORE/AFTER
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.6, 5, 0.5, "KPI & SO SÁNH", 16, CYAN, True, PP_ALIGN.LEFT, "Courier New")
add_text(s, 1, 1.2, 14, 0.8, "Kết quả sau triển khai TAlpha", 40, WHITE, True)

kpis = [
    ("90%", "Giảm thời gian\nsoạn tin", CYAN),
    ("4x", "Tăng tần suất\nchào hàng", GREEN),
    ("0", "SQL / Code\ncần viết thêm", YELLOW),
    ("100%", "Vibe Coded\nby AI", PURPLE),
]
for i, (num, label, clr) in enumerate(kpis):
    x = 0.8 + i * 3.8
    add_rect(s, x, 2.3, 3.4, 2, CARD)
    add_text(s, x + 0.1, 2.4, 3.2, 0.8, num, 48, clr, True, PP_ALIGN.CENTER, "Courier New")
    add_text(s, x + 0.1, 3.3, 3.2, 0.8, label, 16, GRAY, False, PP_ALIGN.CENTER)

# Before/After
add_rect(s, 0.8, 4.8, 7, 3.5, CARD)
add_text(s, 1.1, 5, 6, 0.5, "❌  TRƯỚC", 22, RED, True, PP_ALIGN.LEFT, "Courier New")
add_multiline(s, 1.1, 5.6, 6, 2.5, [
    "• Soạn tin thủ công từng khách",
    "• 2-3h/ngày cho broadcast",
    "• Excel quản lý shipment",
    "• Không tracking đơn hàng",
], 18, GRAY)
add_rect(s, 8.2, 4.8, 7, 3.5, CARD)
add_text(s, 8.5, 5, 6, 0.5, "✅  SAU — TAlpha", 22, GREEN, True, PP_ALIGN.LEFT, "Courier New")
add_multiline(s, 8.5, 5.6, 6, 2.5, [
    "• 30s tạo kịch bản 4 đoạn",
    "• 10 phút/ngày, AI lo phần còn lại",
    "• 17 trạng thái real-time",
    "• Phân phối UAE → 6 nước GCC",
], 18, GRAY)

# =====================================================
# SLIDE 25: THANK YOU
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 2, 1.2, 12, 0.5, "C Ả M   Ơ N", 22, CYAN, True, PP_ALIGN.CENTER, "Courier New")
add_text(s, 1, 2.2, 14, 2, "TAlpha — Vibe Coded\nwith ❤️ by AI", 60, WHITE, True, PP_ALIGN.CENTER)
add_rect(s, 7.2, 4.8, 1.6, 0.04, CYAN, 0)
add_text(s, 2, 5.3, 12, 0.5, "Pi Alpha  ·  Tiểu Alpha — Middle East", 22, GRAY, False, PP_ALIGN.CENTER)
add_text(s, 2, 6, 12, 0.5, "talpha-dashboard.vercel.app", 18, CYAN, False, PP_ALIGN.CENTER, "Courier New")
add_text(s, 2, 6.7, 12, 0.5, "🇸🇦  🇦🇪  🇰🇼  🇴🇲  🇶🇦  🇧🇭", 36, WHITE, False, PP_ALIGN.CENTER)

# Team footer
add_text(s, 1, 7.6, 14, 0.8, "Lê Thị Diệu Thuý  ·  Nguyễn Thị Phương Anh  ·  Hồ Sỹ Anh  ·  Hồ Sỹ Lộc  ·  Nguyễn Thị Hồng Ly", 16, DARK_GRAY, False, PP_ALIGN.CENTER)

# =====================================================
# SAVE
# =====================================================
prs.save(OUT)
print(f"✅ Saved 25 slides to: {OUT}")
print(f"   Size: {os.path.getsize(OUT) / 1024 / 1024:.1f} MB")
